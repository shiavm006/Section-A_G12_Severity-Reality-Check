"""
Severity Reality Check — Capstone 2 Presentation Deck
11 slides following the NST DVA Capstone 2 official PPT outline.
Palette: Midnight Executive (navy + ice blue + coral accent).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── Palette ───────────────────────────────────────────────────────
NAVY     = RGBColor(0x1E, 0x27, 0x61)   # primary
ICE      = RGBColor(0xCA, 0xDC, 0xFC)   # secondary
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CORAL    = RGBColor(0xF9, 0x61, 0x67)   # accent for hero metrics
GRAY     = RGBColor(0x64, 0x74, 0x8B)   # muted text
DARKGRAY = RGBColor(0x1E, 0x29, 0x3B)   # body text
SOFT     = RGBColor(0xF1, 0xF5, 0xF9)   # card background

# ─── Setup ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ─── Helpers ───────────────────────────────────────────────────────

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARKGRAY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARKGRAY,
                bullet_color=NAVY, line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Coloured square bullet
        r1 = p.add_run()
        r1.text = "■  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # Body
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_card(slide, x, y, w, h, fill=WHITE, border=ICE):
    return add_rect(slide, x, y, w, h, fill, border)

def add_footer(slide, page_num, total=11):
    add_rect(slide, 0, 7.2, 13.333, 0.3, NAVY)
    add_text(slide, 0.5, 7.235, 8, 0.25,
             "Severity Reality Check  ·  NST DVA Capstone 2  ·  April 2026",
             size=9, color=ICE, font="Calibri")
    add_text(slide, 12.0, 7.235, 1.0, 0.25,
             f"{page_num} / {total}", size=9, color=ICE, font="Calibri",
             align=PP_ALIGN.RIGHT)

def kpi_card(slide, x, y, w, h, label, value, sublabel="", value_color=CORAL):
    add_card(slide, x, y, w, h, fill=WHITE, border=ICE)
    add_text(slide, x + 0.15, y + 0.15, w - 0.3, 0.35, label,
             size=11, bold=True, color=GRAY)
    add_text(slide, x + 0.15, y + 0.55, w - 0.3, h - 1.1, value,
             size=28, bold=True, color=value_color, font="Calibri")
    if sublabel:
        add_text(slide, x + 0.15, y + h - 0.45, w - 0.3, 0.35, sublabel,
                 size=9, color=GRAY, italic=True)

# ─── Slide 1 — Title ───────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.333, 7.5, NAVY)
# left accent bar
add_rect(s1, 0, 0, 0.5, 7.5, CORAL)

# Eyebrow
add_text(s1, 1.0, 1.6, 11.5, 0.5,
         "NST DVA  ·  CAPSTONE 2  ·  TRANSPORTATION & PUBLIC SAFETY",
         size=12, bold=True, color=ICE, font="Calibri")

# Title
add_text(s1, 1.0, 2.2, 11.5, 1.6,
         "Severity Reality Check",
         size=64, bold=True, color=WHITE, font="Calibri")

# Subtitle
add_text(s1, 1.0, 3.9, 11.5, 0.7,
         "Disentangling true accident severity from congestion,",
         size=22, color=ICE, font="Calibri")
add_text(s1, 1.0, 4.35, 11.5, 0.7,
         "time-of-day, and location bias",
         size=22, color=ICE, font="Calibri")

# Hero stat
add_text(s1, 1.0, 5.4, 6.0, 0.4,
         "ANALYZING 95K US ROAD ACCIDENTS  ·  2016–2023",
         size=11, bold=True, color=CORAL, font="Calibri")

# Footer info
add_rect(s1, 1.0, 6.05, 11.3, 0.025, ICE)
add_text(s1, 1.0, 6.15, 5.8, 0.4,
         "Newton School of Technology",
         size=12, bold=True, color=WHITE, font="Calibri")
add_text(s1, 1.0, 6.5, 5.8, 0.3,
         "Team: Shivam Mittal · Satyam · Members 3–7",
         size=10, color=ICE, font="Calibri")
add_text(s1, 1.0, 6.8, 5.8, 0.3,
         "Mentor: <fill>   ·   Section: <fill>   ·   Team ID: <fill>",
         size=9, color=ICE, italic=True, font="Calibri")

add_text(s1, 7.5, 6.15, 4.8, 0.3,
         "GitHub:  github.com/<team-handle>",
         size=10, color=ICE, font="Calibri")
add_text(s1, 7.5, 6.5, 4.8, 0.3,
         "Tableau:  public.tableau.com/<dashboard>",
         size=10, color=ICE, font="Calibri")
add_text(s1, 7.5, 6.8, 4.8, 0.3,
         "Submission: April 29, 2026",
         size=9, color=ICE, italic=True, font="Calibri")

# ─── Slide 2 — Context & Problem ───────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.333, 7.5, WHITE)
# title bar
add_text(s2, 0.6, 0.4, 12, 0.5, "02  ·  CONTEXT", size=11, bold=True, color=GRAY)
add_text(s2, 0.6, 0.8, 12, 0.9, "Public severity data is widely used —", size=32, bold=True, color=NAVY)
add_text(s2, 0.6, 1.5, 12, 0.9, "but it measures the wrong thing.", size=32, bold=True, color=CORAL)

# Two-column body
add_text(s2, 0.6, 2.7, 6.0, 0.4, "WHAT'S WRONG", size=11, bold=True, color=GRAY)
add_bullets(s2, 0.6, 3.1, 6.0, 3.5, [
    "The most common severity field measures traffic-flow impact (short delay → long delay), not crash physical severity.",
    "A fender-bender at rush hour can register Severity 4. A fatal crash on a quiet road can register Severity 1.",
    "Insurers, DOTs, and logistics planners use this field directly in pricing models, capital plans, and routing — silently absorbing the bias.",
], size=13, line_spacing=1.35)

# Right column — the question
add_card(s2, 7.0, 2.6, 5.8, 4.0, fill=NAVY, border=NAVY)
add_text(s2, 7.25, 2.85, 5.4, 0.4,
         "OUR CORE BUSINESS QUESTION", size=11, bold=True, color=ICE)
add_text(s2, 7.25, 3.3, 5.4, 3.0,
         "How much of what's labeled \"severe\" in public US accident data is genuinely severe — and how should insurers, DOTs, and logistics planners adjust their decisions to use a bias-corrected signal?",
         size=14, color=WHITE)
add_text(s2, 7.25, 6.05, 5.4, 0.4,
         "→  Decision: down-weight raw severity in pricing,",
         size=10, color=CORAL, italic=True)
add_text(s2, 7.25, 6.30, 5.4, 0.4,
         "    capital allocation, and corridor routing.",
         size=10, color=CORAL, italic=True)

add_footer(s2, 2)

# ─── Slide 3 — Data Engineering ────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.333, 7.5, WHITE)
add_text(s3, 0.6, 0.4, 12, 0.5, "03  ·  DATA ENGINEERING", size=11, bold=True, color=GRAY)
add_text(s3, 0.6, 0.8, 12, 0.9, "From 7.7M raw rows to a 95K balanced sample", size=28, bold=True, color=NAVY)

# 4 stat cards
kpi_card(s3, 0.6, 2.2, 2.95, 1.4, "ORIGINAL ROWS", "7.7M", "raw Kaggle dataset", value_color=NAVY)
kpi_card(s3, 3.7, 2.2, 2.95, 1.4, "WORKING SAMPLE", "95,607", "balanced across Sev 1–4", value_color=CORAL)
kpi_card(s3, 6.8, 2.2, 2.95, 1.4, "COLUMNS", "49", "incl. 13 derived features", value_color=NAVY)
kpi_card(s3, 9.9, 2.2, 2.95, 1.4, "STATES COVERED", "49", "6,800+ cities", value_color=NAVY)

# Cleaning steps left
add_text(s3, 0.6, 3.95, 6.0, 0.4, "CLEANING PIPELINE  (notebooks/02_cleaning.ipynb)",
         size=11, bold=True, color=GRAY)
add_bullets(s3, 0.6, 4.35, 6.0, 2.6, [
    "Drop high-null + zero-variance columns (End_Lat/Lng, Country, Turning_Loop)",
    "Validate geographic bounds (lat 24–50, lng −125 to −66)",
    "Cap Distance_mi outliers at 99th percentile",
    "Context-aware imputation (median by State + Month + Hour)",
    "Cast booleans to int for Tableau",
], size=12, line_spacing=1.25)

# Derived features right
add_card(s3, 7.0, 3.95, 5.8, 3.0, fill=SOFT, border=ICE)
add_text(s3, 7.25, 4.10, 5.4, 0.4, "KEY DERIVED FEATURES", size=11, bold=True, color=NAVY)
add_bullets(s3, 7.25, 4.55, 5.4, 2.4, [
    "Severity_Label · Weather_Category · State_Region",
    "Time_of_Day · Season · Is_Rush_Hour · Is_Weekend",
    "Duration_min · Road_Feature_Count",
    "Is True Severe = Sev=4 AND Distance ≥ 0.5 mi",
], size=11, line_spacing=1.3)

add_footer(s3, 3)

# ─── Slide 4 — KPI Framework ───────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.333, 7.5, WHITE)
add_text(s4, 0.6, 0.4, 12, 0.5, "04  ·  KPI FRAMEWORK", size=11, bold=True, color=GRAY)
add_text(s4, 0.6, 0.8, 12, 0.9, "15 decision-relevant KPIs across 3 dashboards", size=28, bold=True, color=NAVY)

# 3 dashboard columns
def kpi_col(x, w, headline, color, sub, items):
    add_card(s4, x, 2.0, w, 4.85, fill=WHITE, border=ICE)
    add_rect(s4, x, 2.0, w, 0.5, color)
    add_text(s4, x + 0.2, 2.05, w - 0.4, 0.4, headline, size=14, bold=True, color=WHITE)
    add_text(s4, x + 0.2, 2.6, w - 0.4, 0.3, sub, size=10, color=GRAY, italic=True)
    add_bullets(s4, x + 0.2, 2.95, w - 0.4, 3.8, items,
                size=11, line_spacing=1.25, bullet_color=color)

kpi_col(0.6, 4.0,
        "D1 — Severity Reality Check",
        NAVY,
        "the hook",
        ["Total Accidents", "Nominal Severe (Sev 4)",
         "True Severe (Sev 4 + Dist ≥ 0.5 mi)",
         "True Severe %  ← headline",
         "Avg Duration of Severe"])

kpi_col(4.75, 4.0,
        "D2 — When & Where",
        CORAL,
        "the diagnosis",
        ["Worst Season",
         "Peak Risk Hour  (12 AM, not rush hour)",
         "Top State by True Severe  (PA, not CA)",
         "Weekend vs Weekday Δ",
         "Cities Tracked"])

kpi_col(8.9, 4.0,
        "D3 — Conditions",
        NAVY,
        "the cause",
        ["Riskiest Weather",
         "Junction Lift  (+0.2–0.5 pts)",
         "Signal Drop  (≈0.4 pts)",
         "Avg Visibility During Severe",
         "Night Severe Share"])

add_footer(s4, 4)

# ─── Slide 5 — Key EDA Insights ────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.333, 7.5, WHITE)
add_text(s5, 0.6, 0.4, 12, 0.5, "05  ·  KEY EDA INSIGHTS", size=11, bold=True, color=GRAY)
add_text(s5, 0.6, 0.8, 12, 0.9, "Five observations that reframe the data", size=28, bold=True, color=NAVY)

# Embed the rush hour plot on the right
plot_path = Path("/sessions/adoring-charming-hawking/mnt/Dva Capstone 2 /reports/eda_plots/severity_rush_hour.png")
if plot_path.exists():
    s5.shapes.add_picture(str(plot_path), Inches(7.6), Inches(2.1),
                          width=Inches(5.3), height=Inches(4.2))
    add_text(s5, 7.6, 6.3, 5.3, 0.3,
             "Severity proportions: Rush Hour vs Non-Rush — barely shift",
             size=9, color=GRAY, italic=True)

# Insights left
add_bullets(s5, 0.6, 2.1, 6.7, 4.6, [
    "Severity-2 dominance dropped after balanced sampling — every tier now has enough N for cross-tier comparison.",
    "Volume peaks at 5 PM rush hour — but per-accident severity is barely different between rush and non-rush.",
    "Weekday volume rises Tue–Wed, falls weekend. Severity does not follow the same curve.",
    "Top-volume states: CA, FL, TX. Top True-Severe state: PA — geography differs once bias is corrected.",
    "Most accidents happen in fair weather. Severe accidents are no exception.",
], size=13, line_spacing=1.4)

add_footer(s5, 5)

# ─── Slide 6 — Advanced Analysis ───────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.333, 7.5, WHITE)
add_text(s6, 0.6, 0.4, 12, 0.5, "06  ·  ADVANCED ANALYSIS", size=11, bold=True, color=GRAY)
add_text(s6, 0.6, 0.8, 12, 0.9, "Statistical tests aligned with the bias hypothesis", size=28, bold=True, color=NAVY)

# 4 test cards
def test_card(x, y, label, method, finding):
    add_card(s6, x, y, 6.0, 1.95, fill=WHITE, border=ICE)
    add_rect(s6, x, y, 0.12, 1.95, NAVY)
    add_text(s6, x + 0.3, y + 0.15, 5.6, 0.4, label, size=12, bold=True, color=NAVY)
    add_text(s6, x + 0.3, y + 0.6, 5.6, 0.4, "Method:  " + method, size=10, color=GRAY, italic=True)
    add_text(s6, x + 0.3, y + 1.0, 5.6, 0.9, finding, size=11, color=DARKGRAY)

test_card(0.6, 2.1,
          "Test 1 — Congestion Bias",
          "Chi-square (Rush Hour × Severity)",
          "Significant by N, but proportions barely shift. Rush-hour volume is NOT the dominant severity driver.")
test_card(6.75, 2.1,
          "Test 2 — Duration × Severity",
          "Kruskal-Wallis H (non-parametric)",
          "Median duration rises monotonically: Sev 1 ≈ 30 min → Sev 4 ≈ 210 min. Validates True-Severe definition.")
test_card(0.6, 4.2,
          "Test 3 — Weather × Severity",
          "Chi-square (top 5 weather × Sev)",
          "Statistically significant but small effect size. Weather is NOT the primary driver of severity.")
test_card(6.75, 4.2,
          "Test 4 — Location Bias",
          "State × Severity cross-tab",
          "PA, WY, MT, NY over-represented in Sev 3+4. Reporting bias is real and must be corrected.")

# Bottom takeaway
add_card(s6, 0.6, 6.4, 12.15, 0.65, fill=NAVY, border=NAVY)
add_text(s6, 0.85, 6.50, 11.7, 0.5,
         "Net finding:  Severity in this dataset is partly real, partly noise. The True-Severe filter (Sev 4 + Dist ≥ 0.5 mi) cleanly separates them.",
         size=12, bold=True, color=WHITE)

add_footer(s6, 6)

# ─── Slide 7 — Tableau Walkthrough ─────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.333, 7.5, WHITE)
add_text(s7, 0.6, 0.4, 12, 0.5, "07  ·  TABLEAU DASHBOARDS", size=11, bold=True, color=GRAY)
add_text(s7, 0.6, 0.8, 12, 0.9, "Three dashboards, one narrative: What → When/Where → Why", size=28, bold=True, color=NAVY)

# 3 columns
def db_col(x, num, name, kpi_strip, charts, color):
    add_card(s7, x, 2.0, 4.0, 4.85, fill=WHITE, border=ICE)
    add_rect(s7, x, 2.0, 4.0, 0.55, color)
    add_text(s7, x + 0.2, 2.05, 3.6, 0.4, f"DASHBOARD {num}", size=10, bold=True, color=WHITE)
    add_text(s7, x + 0.2, 2.32, 3.6, 0.3, name, size=12, bold=True, color=WHITE)
    add_text(s7, x + 0.2, 2.7, 3.6, 0.3, "KPI STRIP", size=9, bold=True, color=GRAY)
    add_text(s7, x + 0.2, 3.0, 3.6, 1.5, kpi_strip, size=10, color=DARKGRAY)
    add_text(s7, x + 0.2, 4.6, 3.6, 0.3, "CHARTS", size=9, bold=True, color=GRAY)
    add_text(s7, x + 0.2, 4.9, 3.6, 1.9, charts, size=10, color=DARKGRAY)

db_col(0.6, 1, "Severity Reality Check",
       "Total · Nominal Severe · True Severe · True Severe % · Avg Severe Duration",
       "Severity Donut · Nominal vs True Severe · Rush Hour Comparison · Hourly Severity Trend",
       NAVY)
db_col(4.75, 2, "When & Where",
       "Worst Season · Peak Risk Hour · Top State by True Severe · Weekend vs Weekday Δ · Cities Tracked",
       "Day × Hour Heatmap · Top 12 States · Time-of-Day Severity Bar · Region Comparison",
       CORAL)
db_col(8.90, 3, "Conditions",
       "Riskiest Weather · Junction Lift · Signal Drop · Avg Visibility During Severe · Night Severe Share",
       "Weather × Severity · Visibility Bucket × Severity · Road Feature Lifts · Temp × Sev by Season",
       NAVY)

# Filters footer
add_text(s7, 0.6, 6.95, 12, 0.3,
         "6+ interactive filters across dashboards: Severity · State · Year · Time-of-Day · Season · Weather · Visibility",
         size=10, color=GRAY, italic=True)

add_footer(s7, 7)

# ─── Slide 8 — Recommendations ─────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, 13.333, 7.5, WHITE)
add_text(s8, 0.6, 0.4, 12, 0.5, "08  ·  RECOMMENDATIONS", size=11, bold=True, color=GRAY)
add_text(s8, 0.6, 0.8, 12, 0.9, "Five actions, each tied to an insight", size=28, bold=True, color=NAVY)

# Header row
add_rect(s8, 0.6, 2.0, 12.1, 0.45, NAVY)
add_text(s8, 0.75, 2.07, 1.0, 0.3, "#",  size=10, bold=True, color=WHITE)
add_text(s8, 1.85, 2.07, 2.6, 0.3, "STAKEHOLDER", size=10, bold=True, color=WHITE)
add_text(s8, 4.55, 2.07, 5.5, 0.3, "RECOMMENDATION", size=10, bold=True, color=WHITE)
add_text(s8, 10.1, 2.07, 2.6, 0.3, "EXPECTED IMPACT", size=10, bold=True, color=WHITE)

recs = [
    ("1", "Insurer", "Discount publicly-reported Severity 4 by ~50%; switch to True-Severe-based pricing input.",
     "Mis-priced urban-policy margin recovered; more accurate risk-tier pricing."),
    ("2", "State DOT", "Prioritize signal installations at top-N junction hotspots (from D3).",
     "≈ 0.4 pts severity reduction per retrofit corridor; measurable per-signal ROI."),
    ("3", "Logistics fleet", "Re-rank corridors by True-Severe Index; shift dispatch off PA + winter night windows.",
     "Lower expected claim frequency and severity per million route-miles."),
    ("4", "Infrastructure planner", "Rebalance away from visibility-only spend; invest in junction redesigns + signal coverage.",
     "2–3× improvement in $/severity-point reduction vs visibility-only."),
    ("5", "Federal grants", "Reweight per-state safety allocation using True-Severe counts, not raw volume.",
     "Capital tracks real severity outcomes, not reporting frequency."),
]
y = 2.55
row_h = 0.85
for i, (n, who, rec, imp) in enumerate(recs):
    bg = SOFT if i % 2 == 0 else WHITE
    add_rect(s8, 0.6, y, 12.1, row_h, bg)
    add_text(s8, 0.75, y + 0.05, 1.0, row_h - 0.1, n,
             size=20, bold=True, color=CORAL, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s8, 1.85, y + 0.10, 2.6, row_h - 0.2, who,
             size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s8, 4.55, y + 0.10, 5.5, row_h - 0.2, rec,
             size=10, color=DARKGRAY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s8, 10.1, y + 0.10, 2.6, row_h - 0.2, imp,
             size=10, color=DARKGRAY, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

add_footer(s8, 8)

# ─── Slide 9 — Impact & Value ──────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, 13.333, 7.5, WHITE)
add_text(s9, 0.6, 0.4, 12, 0.5, "09  ·  IMPACT & VALUE", size=11, bold=True, color=GRAY)
add_text(s9, 0.6, 0.8, 12, 0.9, "Why this work pays off — directional estimates", size=28, bold=True, color=NAVY)

# 3 hero metrics
add_card(s9, 0.6, 2.0, 4.0, 2.5, fill=NAVY, border=NAVY)
add_text(s9, 0.85, 2.20, 3.5, 0.3, "INSURER  ·  REC #1", size=10, bold=True, color=ICE)
add_text(s9, 0.85, 2.6,  3.5, 1.0, "~50%", size=64, bold=True, color=CORAL)
add_text(s9, 0.85, 3.55, 3.5, 0.9, "of Severity-4 in pricing models is over-weighted; bias-correction unlocks ~$100M risk-adj. premium reallocation per $10B book.",
         size=10, color=ICE)

add_card(s9, 4.75, 2.0, 4.0, 2.5, fill=WHITE, border=ICE)
add_text(s9, 5.00, 2.20, 3.5, 0.3, "STATE DOT  ·  REC #2", size=10, bold=True, color=GRAY)
add_text(s9, 5.00, 2.6,  3.5, 1.0, "−0.4 pts", size=54, bold=True, color=NAVY)
add_text(s9, 5.00, 3.55, 3.5, 0.9, "average severity reduction at junctions retrofitted with traffic signals — quantifiable per-installation ROI.",
         size=10, color=DARKGRAY)

add_card(s9, 8.90, 2.0, 4.0, 2.5, fill=WHITE, border=ICE)
add_text(s9, 9.15, 2.20, 3.5, 0.3, "FLEETS  ·  REC #3", size=10, bold=True, color=GRAY)
add_text(s9, 9.15, 2.6,  3.5, 1.0, "5–10%", size=54, bold=True, color=NAVY)
add_text(s9, 9.15, 3.55, 3.5, 0.9, "expected reduction in severe-corridor exposure when routing on True-Severe Index instead of raw volume.",
         size=10, color=DARKGRAY)

# Bottom text
add_text(s9, 0.6, 4.85, 12.1, 0.4, "WHY ACT NOW", size=11, bold=True, color=GRAY)
add_bullets(s9, 0.6, 5.25, 12.1, 1.7, [
    "Real-time accident APIs increasingly feed pricing engines — bias compounds quickly.",
    "Post-COVID vehicle-miles-traveled have rebounded; more decisions flow through the same biased input.",
    "Signal Drop is one of the most cost-effective infrastructure levers — and the data quantifies it.",
], size=12, line_spacing=1.3)

add_footer(s9, 9)

# ─── Slide 10 — Limitations & Next Steps ───────────────────────────
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, 13.333, 7.5, WHITE)
add_text(s10, 0.6, 0.4, 12, 0.5, "10  ·  LIMITATIONS & NEXT STEPS", size=11, bold=True, color=GRAY)
add_text(s10, 0.6, 0.8, 12, 0.9, "Honesty about what this analysis can and cannot say", size=28, bold=True, color=NAVY)

# Two columns
add_text(s10, 0.6, 2.0, 6.0, 0.4, "LIMITATIONS", size=12, bold=True, color=CORAL)
add_bullets(s10, 0.6, 2.45, 6.0, 4.4, [
    "Severity is a flow-impact label, not a crash-injury label — True Severe is a proxy.",
    "0.5-mi threshold is defensible but not absolute; sensitivity tested at 0.25 and 1.0 mi.",
    "All findings are correlational, not causal.",
    "API coverage varies by state; 2020+ volume spike is partly an artifact.",
    "No injury / fatality outcomes available in this dataset.",
    "Balanced sample changes base rates — most KPIs use ratios to control.",
], size=11, line_spacing=1.35)

add_text(s10, 7.0, 2.0, 6.0, 0.4, "NEXT STEPS", size=12, bold=True, color=NAVY)
add_bullets(s10, 7.0, 2.45, 6.0, 4.4, [
    "Join FARS / state crash records to validate True-Severe against real injury data.",
    "Add quasi-experimental causal layer for Signal Drop estimate.",
    "Wrap corrected severity index as a real-time API for fleets and insurers.",
    "Train a gradient-boosted classifier on incoming reports (predict True Severe at intake).",
    "Geospatial drill-down map (Mapbox / Leaflet) for DOT use.",
    "Attach cost-benefit modelling using FHWA retrofit cost data.",
], size=11, line_spacing=1.35, bullet_color=NAVY)

add_footer(s10, 10)

# ─── Slide 11 — Team & Contribution ────────────────────────────────
s11 = prs.slides.add_slide(BLANK)
add_rect(s11, 0, 0, 13.333, 7.5, WHITE)
add_text(s11, 0.6, 0.4, 12, 0.5, "11  ·  TEAM & ARTEFACTS", size=11, bold=True, color=GRAY)
add_text(s11, 0.6, 0.8, 12, 0.9, "Built by Team — every member contributed", size=28, bold=True, color=NAVY)

# Team grid
add_text(s11, 0.6, 2.0, 12, 0.4, "TEAM ROLES", size=11, bold=True, color=GRAY)
roles = [
    ("Project Lead", "Shivam Mittal"),
    ("Data Lead",        "Satyam"),
    ("ETL Lead",         "<Member 3>"),
    ("Analysis Lead",    "<Member 4>"),
    ("Visualization Lead", "<Member 5>"),
    ("Strategy Lead",    "<Member 6>"),
    ("PPT & Quality Lead","<Member 7>"),
]
for i, (role, name) in enumerate(roles):
    col = i % 4
    row = i // 4
    x = 0.6 + col * 3.10
    y = 2.45 + row * 1.05
    add_card(s11, x, y, 2.95, 0.95, fill=SOFT, border=ICE)
    add_text(s11, x + 0.15, y + 0.10, 2.7, 0.35, role,
             size=10, bold=True, color=GRAY)
    add_text(s11, x + 0.15, y + 0.45, 2.7, 0.45, name,
             size=14, bold=True, color=NAVY)

# Artefacts
add_text(s11, 0.6, 4.7, 12, 0.4, "ARTEFACTS  ·  github + tableau", size=11, bold=True, color=GRAY)
add_card(s11, 0.6, 5.15, 12.1, 1.7, fill=NAVY, border=NAVY)
add_text(s11, 0.85, 5.30, 11.5, 0.4, "▸  GitHub Repo:", size=11, bold=True, color=ICE)
add_text(s11, 2.85, 5.30, 9.5, 0.4, "github.com/<team-handle>  (5 notebooks · ETL pipeline · 3 dashboards · report + deck)",
         size=11, color=WHITE)
add_text(s11, 0.85, 5.65, 11.5, 0.4, "▸  Tableau Public:", size=11, bold=True, color=ICE)
add_text(s11, 2.85, 5.65, 9.5, 0.4, "public.tableau.com/<dashboard>  (3 views · 30+ worksheets · 6+ filters)",
         size=11, color=WHITE)
add_text(s11, 0.85, 6.00, 11.5, 0.4, "▸  Project Report:", size=11, bold=True, color=ICE)
add_text(s11, 2.85, 6.00, 9.5, 0.4, "reports/project_report.pdf  (18 sections, ~22 pages)", size=11, color=WHITE)
add_text(s11, 0.85, 6.35, 11.5, 0.4, "▸  This Deck:", size=11, bold=True, color=ICE)
add_text(s11, 2.85, 6.35, 9.5, 0.4, "reports/presentation.pdf  (11 slides)", size=11, color=WHITE)

add_footer(s11, 11)

# ─── Save ──────────────────────────────────────────────────────────
out = Path("/sessions/adoring-charming-hawking/mnt/outputs/deck_build/presentation.pptx")
prs.save(out)
print(f"Saved: {out}  ({out.stat().st_size:,} bytes)")
